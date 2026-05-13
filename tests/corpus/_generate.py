"""Synthetic document corpus generator.

Run this once to populate tests/corpus/*.{docx,pptx,xlsx,pdf}. Outputs are
checked into source so test runs are hermetic.

Usage:
    cd /workspace
    python -m tests.corpus._generate
"""

from __future__ import annotations

import sys
from pathlib import Path

CORPUS_DIR = Path(__file__).parent


def make_small_docx(path: Path) -> None:
    from docx import Document

    doc = Document()
    doc.add_heading("Small DOCX", level=1)
    for i in range(3):
        doc.add_paragraph(f"Paragraph {i + 1} on page-equivalent {i + 1}.")
        doc.add_page_break()
    doc.save(path)


def make_medium_docx(path: Path) -> None:
    from docx import Document

    doc = Document()
    doc.add_heading("Medium DOCX", level=1)
    for i in range(100):
        doc.add_paragraph(f"Page {i + 1}: " + ("lorem ipsum " * 30))
        doc.add_page_break()
    doc.save(path)


def make_simple_pptx(path: Path) -> None:
    from pptx import Presentation

    pres = Presentation()
    layout = pres.slide_layouts[1]
    for i in range(5):
        slide = pres.slides.add_slide(layout)
        slide.shapes.title.text = f"Slide {i + 1}"
        slide.placeholders[1].text = "plain content"
    pres.save(path)


def make_complex_pptx(path: Path) -> None:
    """20 slides with text and shapes (no real images for portability)."""
    from pptx import Presentation
    from pptx.util import Inches

    pres = Presentation()
    layout = pres.slide_layouts[5]  # blank
    for i in range(20):
        slide = pres.slides.add_slide(layout)
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
        tb.text_frame.text = f"complex slide {i + 1}"
        # Decorative rectangle to bulk up the slide
        slide.shapes.add_shape(1, Inches(1), Inches(3), Inches(4), Inches(2))
    pres.save(path)


def make_single_sheet_xlsx(path: Path) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    for col in range(1, 21):
        for row in range(1, 101):
            ws.cell(row=row, column=col, value=f"{row}-{col}")
    wb.save(path)


def make_multi_sheet_xlsx(path: Path) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    wb.remove(wb.active)
    for sheet_i in range(4):
        ws = wb.create_sheet(f"Sheet{sheet_i + 1}")
        rows = 10 * (sheet_i + 1)  # varying sizes
        for r in range(1, rows + 1):
            for c in range(1, 6):
                ws.cell(row=r, column=c, value=f"{r}.{c}")
    wb.save(path)


def make_simple_pdf(path: Path) -> None:
    from reportlab.pdfgen import canvas

    c = canvas.Canvas(str(path))
    for i in range(10):
        c.drawString(100, 750, f"Simple PDF — page {i + 1}")
        c.showPage()
    c.save()


CORPUS = {
    "small.docx": make_small_docx,
    "medium.docx": make_medium_docx,
    "simple.pptx": make_simple_pptx,
    "complex.pptx": make_complex_pptx,
    "single_sheet.xlsx": make_single_sheet_xlsx,
    "multi_sheet.xlsx": make_multi_sheet_xlsx,
    "simple.pdf": make_simple_pdf,
}


def main() -> int:
    for filename, maker in CORPUS.items():
        path = CORPUS_DIR / filename
        if path.exists():
            print(f"  exists: {filename}")
            continue
        print(f"generate: {filename}")
        try:
            maker(path)
        except ImportError as e:
            print(f"  skip {filename}: missing dep ({e.name})", file=sys.stderr)
    return 0


if __name__ == "__main__":
    sys.exit(main())
